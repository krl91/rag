"""
Tests du module d'export MS Office (export/office_writer.py).

Écriture réelle sur disque (tmp_path) puis relecture avec les mêmes
bibliothèques que le reste du projet (python-docx, python-pptx, openpyxl) —
aucun appel réseau ni LLM.
"""

from pathlib import Path

from kb_smart_metering.export.office_writer import (
    markdown_to_docx,
    markdown_to_pptx,
    rows_to_xlsx,
)

_SAMPLE_MD = """# Titre principal

## Résumé

Un paragraphe de résumé.

## Décisions

- Première décision
- Deuxième décision

## Sources

| Source | URL |
|---|---|
| JIRA-1 | https://example.com/1 |
| JIRA-2 | https://example.com/2 |
"""


class TestMarkdownToDocx:
    def test_titres_et_paragraphe(self, tmp_path: Path) -> None:
        import docx

        out = tmp_path / "out.docx"
        markdown_to_docx(_SAMPLE_MD, out)

        doc = docx.Document(str(out))
        styles_and_text = [(p.style.name, p.text) for p in doc.paragraphs if p.text.strip()]

        assert ("Heading 1", "Titre principal") in styles_and_text
        assert ("Heading 2", "Résumé") in styles_and_text
        assert ("Normal", "Un paragraphe de résumé.") in styles_and_text

    def test_liste_a_puces(self, tmp_path: Path) -> None:
        import docx

        out = tmp_path / "out.docx"
        markdown_to_docx(_SAMPLE_MD, out)

        doc = docx.Document(str(out))
        bullets = [p.text for p in doc.paragraphs if p.style.name == "List Bullet"]
        assert "Première décision" in bullets
        assert "Deuxième décision" in bullets

    def test_tableau(self, tmp_path: Path) -> None:
        import docx

        out = tmp_path / "out.docx"
        markdown_to_docx(_SAMPLE_MD, out)

        doc = docx.Document(str(out))
        assert len(doc.tables) == 1
        rows = [[c.text for c in row.cells] for row in doc.tables[0].rows]
        assert rows[0] == ["Source", "URL"]
        assert rows[1] == ["JIRA-1", "https://example.com/1"]

    def test_cree_repertoire_parent(self, tmp_path: Path) -> None:
        out = tmp_path / "sous" / "dossier" / "out.docx"
        markdown_to_docx("# Titre", out)
        assert out.exists()

    def test_markdown_vide(self, tmp_path: Path) -> None:
        out = tmp_path / "vide.docx"
        markdown_to_docx("", out)
        assert out.exists()


class TestMarkdownToPptx:
    def test_titre_explicite_une_seule_diapositive_titre(self, tmp_path: Path) -> None:
        from pptx import Presentation

        out = tmp_path / "out.pptx"
        markdown_to_pptx(_SAMPLE_MD, out, title="Titre principal")

        prs = Presentation(str(out))
        title_texts = [
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text == "Titre principal"
        ]
        # Une seule diapositive doit porter ce texte (pas de doublon H1 + --title)
        assert len(title_texts) == 1

    def test_h2_cree_une_diapositive_par_section(self, tmp_path: Path) -> None:
        from pptx import Presentation

        out = tmp_path / "out.pptx"
        markdown_to_pptx(_SAMPLE_MD, out, title="Titre principal")

        prs = Presentation(str(out))
        # 1 diapositive de titre + Résumé + Décisions + Sources = 4
        assert len(prs.slides) == 4

    def test_puces_dans_le_corps(self, tmp_path: Path) -> None:
        from pptx import Presentation

        out = tmp_path / "out.pptx"
        markdown_to_pptx(_SAMPLE_MD, out, title="Titre principal")

        prs = Presentation(str(out))
        all_text = "\n".join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "Première décision" in all_text
        assert "Deuxième décision" in all_text

    def test_sans_titre_explicite(self, tmp_path: Path) -> None:
        """Sans --title, le H1 du markdown sert de diapositive de titre."""
        from pptx import Presentation

        out = tmp_path / "out.pptx"
        markdown_to_pptx(_SAMPLE_MD, out, title=None)

        prs = Presentation(str(out))
        assert len(prs.slides) == 4


class TestRowsToXlsx:
    def test_lignes_et_en_tetes(self, tmp_path: Path) -> None:
        import openpyxl

        out = tmp_path / "out.xlsx"
        rows_to_xlsx([{"id": 1, "nom": "A"}, {"id": 2, "nom": "B"}], out)

        wb = openpyxl.load_workbook(str(out))
        ws = wb.active
        values = list(ws.iter_rows(values_only=True))
        assert values[0] == ("id", "nom")
        assert values[1] == (1, "A")
        assert values[2] == (2, "B")

    def test_cles_heterogenes_union_des_colonnes(self, tmp_path: Path) -> None:
        import openpyxl

        out = tmp_path / "out.xlsx"
        rows_to_xlsx([{"a": 1}, {"a": 2, "b": 3}], out)

        wb = openpyxl.load_workbook(str(out))
        ws = wb.active
        values = list(ws.iter_rows(values_only=True))
        assert values[0] == ("a", "b")
        assert values[1] == (1, None)
        assert values[2] == (2, 3)

    def test_valeur_liste_serialisee_en_json(self, tmp_path: Path) -> None:
        import openpyxl

        out = tmp_path / "out.xlsx"
        rows_to_xlsx([{"tags": ["x", "y"]}], out)

        wb = openpyxl.load_workbook(str(out))
        ws = wb.active
        values = list(ws.iter_rows(values_only=True))
        assert values[1][0] == '["x", "y"]'

    def test_liste_vide(self, tmp_path: Path) -> None:
        out = tmp_path / "vide.xlsx"
        rows_to_xlsx([], out)
        assert out.exists()

    def test_nom_feuille_tronque_31_caracteres(self, tmp_path: Path) -> None:
        import openpyxl

        out = tmp_path / "out.xlsx"
        long_name = "x" * 50
        rows_to_xlsx([{"a": 1}], out, sheet_name=long_name)

        wb = openpyxl.load_workbook(str(out))
        assert len(wb.active.title) <= 31
