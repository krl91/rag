"""
Export MS Office (docx, pptx, xlsx) — écriture symétrique à la lecture déjà
assurée par extractors/office.py.

docx/pptx sont produits à partir du Markdown déjà généré par ce projet
(rendu de `kb ask`/`kb search`, sortie de `kb docgen`) via un parseur
Markdown volontairement minimal — pas un moteur CommonMark complet, conçu
pour le Markdown que ce projet génère lui-même (titres #/##/###, listes à
puces "- "/"* ", tableaux "| a | b |").

xlsx est produit à partir de données tabulaires (liste de dict), typiquement
un export JSON déjà produit par `kb ask --json` ou toute liste plate.

Aucun appel LLM ici.
"""

import json
import logging
import re
from pathlib import Path
from typing import Any, Optional

logger = logging.getLogger(__name__)

_SEPARATOR_RE = re.compile(r"^\|?[\s:|-]+\|?$")


class _Block:
    """Un bloc Markdown identifié par le parseur minimal ci-dessous."""

    __slots__ = ("kind", "text", "level", "rows")

    def __init__(
        self,
        kind: str,
        text: str = "",
        level: int = 0,
        rows: Optional[list[list[str]]] = None,
    ) -> None:
        self.kind = kind  # "heading" | "bullet" | "paragraph" | "table"
        self.text = text
        self.level = level
        self.rows = rows or []


def _parse_markdown_blocks(markdown_text: str) -> list[_Block]:
    """Découpe un texte Markdown (déjà produit par ce projet) en blocs simples."""
    blocks: list[_Block] = []
    lines = markdown_text.splitlines()
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()

        if not stripped or stripped == "---":
            i += 1
            continue

        heading_match = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if heading_match:
            blocks.append(
                _Block("heading", heading_match.group(2).strip(), len(heading_match.group(1)))
            )
            i += 1
            continue

        if stripped.startswith(("- ", "* ")):
            blocks.append(_Block("bullet", stripped[2:].strip()))
            i += 1
            continue

        if stripped.startswith("|"):
            table_lines: list[str] = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i].strip())
                i += 1
            rows = [
                [cell.strip() for cell in row.strip("|").split("|")]
                for row in table_lines
                if not _SEPARATOR_RE.fullmatch(row)
            ]
            if rows:
                blocks.append(_Block("table", rows=rows))
            continue

        blocks.append(_Block("paragraph", stripped))
        i += 1

    return blocks


# ---------------------------------------------------------------------------
# Word (.docx)
# ---------------------------------------------------------------------------


def markdown_to_docx(markdown_text: str, out_path: Path) -> None:
    """Convertit un texte Markdown (déjà produit par ce projet) en document Word."""
    import docx

    blocks = _parse_markdown_blocks(markdown_text)
    doc = docx.Document()

    for block in blocks:
        if block.kind == "heading":
            doc.add_heading(block.text, level=min(max(block.level, 1), 9))
        elif block.kind == "bullet":
            doc.add_paragraph(block.text, style="List Bullet")
        elif block.kind == "table" and block.rows:
            n_cols = max(len(r) for r in block.rows)
            table = doc.add_table(rows=0, cols=n_cols)
            table.style = "Table Grid"
            for row in block.rows:
                cells = table.add_row().cells
                for idx, val in enumerate(row):
                    if idx < n_cols:
                        cells[idx].text = val
        elif block.kind == "paragraph":
            doc.add_paragraph(block.text)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))
    logger.info("docx écrit : %s (%d bloc(s))", out_path, len(blocks))


# ---------------------------------------------------------------------------
# PowerPoint (.pptx)
# ---------------------------------------------------------------------------


def markdown_to_pptx(markdown_text: str, out_path: Path, title: Optional[str] = None) -> None:
    """
    Convertit un texte Markdown en présentation PowerPoint.

    Chaque titre de niveau 1 ou 2 démarre une nouvelle diapositive ; les
    puces/paragraphes qui suivent deviennent le corps de la diapositive.
    """
    from pptx import Presentation

    blocks = _parse_markdown_blocks(markdown_text)
    prs = Presentation()
    content_layout = prs.slide_layouts[1]

    title_used = False

    def ensure_title_slide(text: str) -> None:
        nonlocal title_used
        if title_used:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = text
        title_used = True

    if title:
        ensure_title_slide(title)

    body_tf = None
    body_empty = True

    def start_slide(slide_title: str) -> None:
        nonlocal body_tf, body_empty
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_title
        body_tf = slide.placeholders[1].text_frame
        body_tf.clear()
        body_empty = True

    def add_line(text: str) -> None:
        nonlocal body_empty
        if body_tf is None:
            start_slide(title or "Diapositive")
        assert body_tf is not None
        if body_empty:
            paragraph = body_tf.paragraphs[0]
            body_empty = False
        else:
            paragraph = body_tf.add_paragraph()
        paragraph.text = text

    # Titre de niveau 1 → diapositive de titre (une seule, la première
    # rencontrée si `title` n'était pas déjà fourni). Titre de niveau 2 →
    # nouvelle diapositive de contenu. Niveaux supérieurs → simple ligne.
    for block in blocks:
        if block.kind == "heading" and block.level == 1:
            ensure_title_slide(block.text)
        elif block.kind == "heading" and block.level == 2:
            start_slide(block.text)
        elif block.kind == "bullet":
            add_line(block.text)
        elif block.kind == "heading":
            add_line(block.text)
        elif block.kind == "paragraph":
            add_line(block.text)
        elif block.kind == "table" and block.rows:
            add_line(" | ".join(block.rows[0]))
            for row in block.rows[1:]:
                add_line(" | ".join(row))

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    logger.info("pptx écrit : %s (%d diapositive(s))", out_path, len(prs.slides))


# ---------------------------------------------------------------------------
# Excel (.xlsx)
# ---------------------------------------------------------------------------


def rows_to_xlsx(rows: list[dict[str, Any]], out_path: Path, sheet_name: str = "Feuille1") -> None:
    """Écrit une liste de dictionnaires (lignes plates) dans un classeur Excel."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet_name[:31] or "Feuille1"  # limite Excel : 31 caractères

    if rows:
        headers = list(dict.fromkeys(key for row in rows for key in row.keys()))
        ws.append(headers)
        for row in rows:
            ws.append([_stringify(row.get(h)) for h in headers])

    out_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(str(out_path))
    logger.info("xlsx écrit : %s (%d ligne(s))", out_path, len(rows))


def _stringify(value: Any) -> Any:
    """Rend une valeur compatible cellule Excel (les listes/dict deviennent du JSON)."""
    if isinstance(value, (list, dict)):
        return json.dumps(value, ensure_ascii=False)
    return value
